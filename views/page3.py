import streamlit as st
import re
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from fpdf import FPDF
import speech_recognition as sr
from groq import Groq
from dotenv import load_dotenv
import os

# Load Groq API key from .env
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API")
if not GROQ_API_KEY:
    st.error("GROQ_API not set in .env file")
    st.stop()
groq_client = Groq(api_key=GROQ_API_KEY)

# --- Notes Generator Prompt ---
NOTES_GENERATOR_PROMPT = """
You are an AI-powered study buddy specializing in generating detailed, exam-oriented notes from uploaded PDF documents or provided syllabus content. Your goal is to provide a comprehensive, well-structured, and easy-to-understand summary covering all main topics and sub-topics, tailored for undergraduate students. Do not simply copy text; instead, synthesize, rephrase, and explain all key points in your own words.

---

### Instructions for Generating Notes:

1. *Comprehensive Coverage:*
   * Read the entire provided document or syllabus and identify all main topics and sub-topics.
   * For each topic and sub-topic, provide a detailed, exam-focused explanation, including definitions, key concepts, important terms, methods, algorithms, and essential points.

2. *Clarity & Structure:*
   * Use clear, simple language and illustrative examples where helpful.
   * Organize the notes with clear headings for each topic and sub-topic.
   * Use bullet points or numbered lists to enhance clarity and aid quick revision.

3. *Visuals & Diagrams:*
   * For any topic that typically requires a diagram or is present with one in the PDF, clearly state the topic name and instruct the user to "refer to the PDF for relevant diagrams."
   * Do not attempt to describe complex diagrams textually.

4. *Exam Focus:*
   * Focus strictly on key concepts and essential information likely to be tested in an exam, avoiding irrelevant details.
   * Ensure the notes are detailed enough for an undergraduate exam, but as concise as possible without losing important information.

5. *Output:*
   * Generate the entire set of notes in one go, covering all topics and sub-topics in the document. Do not pause or ask for user input.
   * The output should be suitable for direct use as exam revision notes.
"""

# --- File extraction ---
def extract_pdf(file):
    # Return a list of page texts, but also debug log the length and preview
    pages = []
    try:
        # Save uploaded file to disk for PyPDF2 compatibility
        import tempfile
        file.seek(0)
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
            tmp.write(file.read())
            tmp.flush()
            tmp_path = tmp.name
        file.seek(0)
        reader = PdfReader(tmp_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            st.info(f"Page {i+1} extracted text length: {len(text) if text else 0}")
            if text:
                pages.append(text)
            else:
                pages.append('')
        if not pages or all(not p.strip() for p in pages):
            st.warning(f"PDF extraction: All pages empty. (Total pages: {len(reader.pages)})")
        else:
            st.info(f"PDF extraction: {len(pages)} pages, first page preview: {repr(pages[0][:100]) if pages[0] else 'EMPTY'}")
    except Exception as e:
        st.error(f"PDF extraction error: {e}")
    return pages

def extract_docx(file):
    return [p.text for p in docx.Document(file).paragraphs if p.text.strip()]

def extract_pptx(file):
    return [shape.text.strip() for slide in Presentation(file).slides
            for shape in slide.shapes
            if hasattr(shape, 'text') and shape.text.strip()]

# --- AI-powered Document Chat and Notes using Groq ---
class DocumentUnderstanding:
    def __init__(self):
        self.full_text = ""
        self.notes_pointer = 0
        self.notes_chunks = []
        self.last_notes_output = ""

    def save_to_session(self):
        st.session_state['docu_full_text'] = self.full_text
        st.session_state['docu_notes_pointer'] = self.notes_pointer
        st.session_state['docu_notes_chunks'] = self.notes_chunks
        st.session_state['docu_last_notes_output'] = self.last_notes_output

    def load_from_session(self):
        self.full_text = st.session_state.get('docu_full_text', "")
        self.notes_pointer = st.session_state.get('docu_notes_pointer', 0)
        self.notes_chunks = st.session_state.get('docu_notes_chunks', [])
        self.last_notes_output = st.session_state.get('docu_last_notes_output', "")

    @property
    def notes_confirmed(self):
        # Use session state to persist confirmation across reruns
        return st.session_state.get('notes_confirmed', False)

    @notes_confirmed.setter
    def notes_confirmed(self, value):
        st.session_state['notes_confirmed'] = value

    def process_document(self, pages):
        # Join with newlines to preserve paragraph/topic structure for splitting
        # Also debug log the number of non-empty pages and preview
        non_empty_pages = [p for p in pages if p and p.strip()]
        self.full_text = "\n".join(non_empty_pages)
        st.info(f"process_document: {len(non_empty_pages)} non-empty pages. Preview: {repr(self.full_text[:100])}")
        self.notes_pointer = 0
        self.notes_chunks = self._split_topics_for_notes()
        self.notes_confirmed = False
        self.last_notes_output = ""
        self.save_to_session()

    def _split_topics_for_notes(self):
        # No chunking needed, just return the whole document as one topic
        return [["Entire Document"]]

    def generate_notes(self, confirm=False, continue_notes=False):
        # Always reload state from session to persist across reruns
        self.load_from_session()
        # Debug: show full_text length and preview
        if not self.full_text or not self.full_text.strip():
            return "No document content available to generate notes. (Debug: full_text is empty)"
        if len(self.full_text) < 30:
            return f"No document content available to generate notes. (Debug: full_text too short: {repr(self.full_text[:100])})"
        # No confirmation or chunking, just generate all notes in one go
        prompt = NOTES_GENERATOR_PROMPT + "\n\nDocument Content:\n" + self.full_text[:12000]
        try:
            chat_completion = groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama-3.1-8b-instant",
            )
            self.last_notes_output = chat_completion.choices[0].message.content.strip()
            self.save_to_session()
            return self.last_notes_output
        except Exception as e:
            st.error(f"Groq API error: {e}")
            return "Could not generate notes. Please try again."

    def answer_question(self, question):
        if not self.full_text.strip():
            return "No document content available to answer questions."
        try:
            prompt = (
                "You are an expert assistant. Answer the user's question using ONLY the following document content. "
                "If the answer is not present, say so.\n\n"
                f"Document Content:\n{self.full_text[:12000]}\n\nQuestion: {question}\nAnswer:"
            )
            chat_completion = groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama-3.1-8b-instant",
            )
            return chat_completion.choices[0].message.content.strip()
        except Exception as e:
            st.error(f"Groq API error: {e}")
            return "Could not generate an answer. Please try a different question."

# --- Voice input ---
def listen():
    r = sr.Recognizer()
    try:
        with sr.Microphone() as source:
            st.toast("Listening... Speak now (5 second timeout)")
            r.adjust_for_ambient_noise(source)
            audio = r.listen(source, timeout=5, phrase_time_limit=5)
        text = r.recognize_google(audio)
        return text
    except Exception as e:
        st.toast(f"Couldn't understand speech: {e}")
        return None

# --- Streamlit UI ---
st.title(":material/robot_2: AI Buddy")

if "library" not in st.session_state:
    st.session_state.library = {}
if "messages" not in st.session_state:
    st.session_state.messages = []
if "document_understanding" not in st.session_state or not isinstance(st.session_state.document_understanding, DocumentUnderstanding):
    st.session_state.document_understanding = DocumentUnderstanding()

# Tabs: Notes Generator and Chat
tab1, tab2 = st.tabs(["Notes Generator", "Chat"])
with tab1:
    # --- Unified Smart Notes: Either Camera or File Upload, but not both ---
    st.subheader(":material/assignment_late: AI Notes Generator")
    input_method = st.radio(
        "Choose input method for notes generation:",
        ("Upload PDF/DOCX/PPTX", "Take Board Photo"),
        horizontal=True
    )
    extracted_text = ""
    uploaded_file = None
    picture = None
    if input_method == "Upload PDF/DOCX/PPTX":
        uploaded_file = st.file_uploader("Upload PDF/DOCX/PPTX", type=["pdf", "docx", "pptx"])
        if uploaded_file:
            if uploaded_file.name not in st.session_state.library:
                with st.spinner(":material/text_compare: Analyzing document..."):
                    try:
                        if uploaded_file.type == "application/pdf":
                            pages = extract_pdf(uploaded_file)
                        elif uploaded_file.type.endswith("document"):
                            pages = extract_docx(uploaded_file)
                        elif uploaded_file.type.endswith("presentation"):
                            pages = extract_pptx(uploaded_file)
                        else:
                            st.error("Unsupported file type")
                            st.stop()
                        st.session_state.document_understanding.process_document(pages)
                        st.session_state.library[uploaded_file.name] = pages
                        extracted_text = st.session_state.document_understanding.full_text
                    except Exception as e:
                        st.error(f"Error processing document: {e}")
                        st.stop()
            else:
                # Already processed, just load text
                extracted_text = "\n".join(st.session_state.library[uploaded_file.name])
    elif input_method == "Take Board Photo":
        enable_camera = st.checkbox(":material/camera: Enable camera")
        picture = st.camera_input("Take a picture of the classroom board", disabled=not enable_camera)
        if picture:
            st.image(picture, caption="Board Photo", use_container_width=True)
            import tempfile
            temp_img_path = None
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_img:
                tmp_img.write(picture.getvalue())
                temp_img_path = tmp_img.name
            try:
                from utils.plag.vision import extract_text_from_image
                extracted_text = extract_text_from_image(temp_img_path)
                st.success("Text extracted from board photo!")
                st.text_area("Extracted Text", extracted_text, height=200)
            except Exception as e:
                st.error(f"Vision model error: {e}")

    # Only show notes generation if we have extracted text from either method
    if extracted_text and len(extracted_text.strip()) > 10:
        docu = st.session_state.document_understanding
        docu.full_text = extracted_text
        docu.save_to_session()
        if st.button("Generate Notes!"):
            notes_output = docu.generate_notes()
            st.session_state['last_notes_output'] = notes_output
            st.markdown(notes_output)
        elif 'last_notes_output' in st.session_state:
            st.markdown(st.session_state['last_notes_output'])

        # PDF Download button for generated notes (with Unicode support using PyFPDF + DejaVu font)
        if 'last_notes_output' in st.session_state and st.session_state['last_notes_output'].strip():
            from fpdf import FPDF
            import tempfile
            import os
            class PDF(FPDF):
                def header(self):
                    self.set_font('DejaVu', 'B', 14)
                    self.cell(0, 10, 'AI-Generated Study Notes', ln=True, align='C')
                    self.ln(5)
            pdf = PDF()
            # Find DejaVuSans.ttf font (try system, matplotlib, or fallback to Arial if not found)
            font_path = None
            try:
                from matplotlib import font_manager
                font_path = font_manager.findfont("DejaVu Sans")
            except Exception:
                # Fallback to common Linux path
                if os.path.exists("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"):
                    font_path = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
            if font_path and os.path.exists(font_path):
                pdf.add_font('DejaVu', '', font_path, uni=True)
                pdf.add_font('DejaVu', 'B', font_path, uni=True)
                pdf.set_font('DejaVu', '', 12)
            else:
                pdf.set_font('Arial', '', 12)
            pdf.add_page()
            pdf.set_auto_page_break(auto=True, margin=15)
            # Split notes into lines and add to PDF, handling long lines
            for line in st.session_state['last_notes_output'].split('\n'):
                pdf.multi_cell(0, 10, line)
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                pdf.output(tmp_pdf.name)
                tmp_pdf.seek(0)
                st.download_button(
                    label="Download Notes as PDF",
                    data=tmp_pdf.read(),
                    file_name="AI_Study_Notes.pdf",
                    mime="application/pdf"
                )
with tab2:
    # Always reload document state for chat tab as well
    st.session_state.document_understanding.load_from_session()
    if st.session_state.library:
        doc_names = list(st.session_state.library.keys())
    col1, col2 = st.columns([0.92, 0.08], vertical_alignment="bottom", gap="small")
    with col2:
        if st.button(":material/mic:"):
            query = listen()
            if query:
                st.session_state.last_question = query
    with col1:
        query_text = st.text_input(label="", placeholder="Ask Anything", value=st.session_state.get('last_question', ''))
    if query_text:
        st.session_state.messages.append({"role": "user", "content": query_text})
        # Ensure document_understanding is loaded with latest session state
        st.session_state.document_understanding.load_from_session()
        answer = st.session_state.document_understanding.answer_question(query_text)
        formatted_answer = f"From the document:\n\n{answer}"
        st.session_state.messages.append({"role": "assistant", "content": formatted_answer})
    with st.container(border=True):
        for message in reversed(st.session_state.messages[-5:]):
            with st.chat_message(message["role"]):
                st.markdown(message["content"])